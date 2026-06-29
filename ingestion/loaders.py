"""Multi-format document loading.

Routes each supported file type to the cleanest available loader and returns a
list of LangChain Documents with format-appropriate metadata, normalized to
clean text / Markdown before chunking.

Engines:
  - PDF              -> PyMuPDF (page numbers preserved; see ingestion.pdf_loader)
  - PPTX             -> MarkItDown (one Document per slide via slide markers),
                        falling back to python-pptx if markers are absent
  - XLSX             -> pandas + openpyxl (per sheet, row blocks, column headers)
  - CSV              -> pandas (row blocks, column headers)
  - DOCX/TXT/MD/JSON -> MarkItDown (clean Markdown), plain-text read as fallback

A failure to read/convert raises DocumentLoadError with a user-friendly message.
"""

import re

from langchain_core.documents import Document

from config.settings import ROWS_PER_CHUNK
from ingestion.pdf_loader import load_pdf

# Formats converted to free text/Markdown (no per-unit metadata).
TEXT_LIKE_FORMATS = ("docx", "txt", "md", "json")

_SLIDE_MARKER = re.compile(r"<!--\s*Slide number:\s*(\d+)\s*-->", re.IGNORECASE)

_markitdown = None


class DocumentLoadError(Exception):
    """Raised when a document cannot be read or converted."""


# --- helpers ---------------------------------------------------------------

def _get_markitdown():
    global _markitdown
    if _markitdown is None:
        from markitdown import MarkItDown
        _markitdown = MarkItDown()
    return _markitdown

def _markitdown_raw(file_path: str) -> str:
    """Convert a file to Markdown text via MarkItDown (no post-cleaning)."""
    result = _get_markitdown().convert(file_path)
    return result.text_content or ""


def _clean_text(text: str) -> str:
    """Light normalization that preserves Markdown/line structure.

    Unlike the aggressive PDF whitespace collapse, this keeps newlines (so tables
    and lists survive) while trimming trailing spaces and runs of blank lines.
    """
    if not text:
        return ""
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    lines = [line.rstrip() for line in text.split("\n")]
    cleaned = "\n".join(lines)
    cleaned = re.sub(r"\n{3,}", "\n\n", cleaned)
    return cleaned.strip()


def _read_text_file(file_path: str) -> str:
    with open(file_path, "r", encoding="utf-8", errors="replace") as fh:
        return fh.read()


# --- per-format loaders ----------------------------------------------------

def _load_text_like(file_path: str, file_type: str) -> list[Document]:
    """DOCX/TXT/MD/JSON via MarkItDown, with a plain-read fallback for text."""
    text = ""
    try:
        text = _clean_text(_markitdown_raw(file_path))
    except Exception as e:
        if file_type in ("txt", "md", "json"):
            text = _clean_text(_read_text_file(file_path))
        else:
            raise DocumentLoadError(
                f"Could not convert the {file_type.upper()} document."
            ) from e
    if not text.strip() and file_type in ("txt", "md", "json"):
        text = _clean_text(_read_text_file(file_path))
    if not text.strip():
        raise DocumentLoadError(f"No readable text found in the {file_type.upper()} file.")
    return [Document(page_content=text, metadata={})]


def _load_pptx_python(file_path: str) -> list[Document]:
    """One Document per slide using python-pptx, with slide_number metadata."""
    from pptx import Presentation

    prs = Presentation(file_path)
    docs = []
    for index, slide in enumerate(prs.slides, start=1):
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs)
                    if line.strip():
                        parts.append(line)
            if shape.has_table:
                for row in shape.table.rows:
                    parts.append(" | ".join(cell.text for cell in row.cells))
        content = _clean_text("\n".join(parts))
        if content:
            docs.append(Document(page_content=content, metadata={"slide_number": index}))
    if not docs:
        raise DocumentLoadError("No readable text found in the presentation.")
    return docs


def _split_pptx_markdown(text: str) -> list[Document]:
    """Split MarkItDown PPTX output into per-slide Documents using slide markers."""
    matches = list(_SLIDE_MARKER.finditer(text or ""))
    if not matches:
        return []
    docs = []
    for i, match in enumerate(matches):
        slide_no = int(match.group(1))
        start = match.end()
        end = matches[i + 1].start() if i + 1 < len(matches) else len(text)
        content = _clean_text(text[start:end])
        if content:
            docs.append(Document(page_content=content, metadata={"slide_number": slide_no}))
    return docs


def _load_pptx(file_path: str) -> list[Document]:
    """PPTX via MarkItDown (primary); fall back to python-pptx."""
    try:
        docs = _split_pptx_markdown(_markitdown_raw(file_path))
        if docs:
            return docs
    except Exception:
        pass
    return _load_pptx_python(file_path)


def _dataframe_to_docs(df, sheet_name: str | None) -> list[Document]:
    """Render a dataframe into row-block Documents, repeating the header row.

    Each chunk is a small Markdown-ish table preserving column names, with
    1-indexed row_start/row_end (and sheet_name for spreadsheets) in metadata.
    """
    headers = [str(h) for h in df.columns]
    header_line = " | ".join(headers)
    separator = " | ".join("---" for _ in headers)
    docs = []
    total = len(df)
    for start in range(0, total, ROWS_PER_CHUNK):
        block = df.iloc[start:start + ROWS_PER_CHUNK]
        lines = [header_line, separator]
        for _, row in block.iterrows():
            lines.append(" | ".join(str(v) for v in row.tolist()))
        meta = {"row_start": start + 1, "row_end": start + len(block)}
        if sheet_name is not None:
            meta["sheet_name"] = sheet_name
        docs.append(Document(page_content="\n".join(lines), metadata=meta))
    return docs


def _load_xlsx(file_path: str) -> list[Document]:
    import pandas as pd

    sheets = pd.read_excel(file_path, sheet_name=None, dtype=str)
    docs = []
    for sheet_name, df in sheets.items():
        if df.empty:
            continue
        docs.extend(_dataframe_to_docs(df.fillna(""), sheet_name=str(sheet_name)))
    if not docs:
        raise DocumentLoadError("No data found in the spreadsheet.")
    return docs


def _load_csv(file_path: str) -> list[Document]:
    import pandas as pd

    df = pd.read_csv(file_path, dtype=str).fillna("")
    docs = _dataframe_to_docs(df, sheet_name=None)
    if not docs:
        raise DocumentLoadError("No rows found in the CSV file.")
    return docs


# --- dispatcher ------------------------------------------------------------

def load_document(file_path: str, file_type: str) -> list[Document]:
    """Load a file of the given type into a list of normalized Documents.

    Raises DocumentLoadError (with a user-friendly message) on any read/convert
    failure or unsupported type.
    """
    try:
        if file_type == "pdf":
            return load_pdf(file_path)
        if file_type == "pptx":
            return _load_pptx(file_path)
        if file_type == "xlsx":
            return _load_xlsx(file_path)
        if file_type == "csv":
            return _load_csv(file_path)
        if file_type in TEXT_LIKE_FORMATS:
            return _load_text_like(file_path, file_type)
    except DocumentLoadError:
        raise
    except Exception as e:
        raise DocumentLoadError(
            f"Could not read the {file_type.upper()} file. It may be corrupted or unsupported."
        ) from e

    raise DocumentLoadError(f"Unsupported file type: {file_type}.")
